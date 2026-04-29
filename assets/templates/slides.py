#!/usr/bin/env python3
"""
gen_slides.py - parchment design system slide deck generator

用法：
  pip install python-pptx --break-system-packages
  python3 gen_slides.py

输出：
  output.pptx (16:9 宽屏, parchment 风格)

这是一个模板脚本 - 填充自己的内容后直接运行。
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ═══════════════════════════════════════════════════════════
# Design System Constants
# ═══════════════════════════════════════════════════════════

# 色板
PARCHMENT   = RGBColor(0xF6, 0xF0, 0xEA)
IVORY       = RGBColor(0xFB, 0xF7, 0xF3)
BRAND       = RGBColor(0xB8, 0x3D, 0x2E)
BRAND_DEEP  = RGBColor(0x7E, 0x24, 0x1C)
NEAR_BLACK  = RGBColor(0x19, 0x15, 0x14)
DARK_WARM   = RGBColor(0x3A, 0x31, 0x2D)
CHARCOAL    = RGBColor(0x4B, 0x3E, 0x39)
OLIVE       = RGBColor(0x5A, 0x4A, 0x43)
STONE       = RGBColor(0x85, 0x77, 0x6F)
BORDER      = RGBColor(0xE9, 0xDE, 0xD4)
WHITE       = RGBColor(0xff, 0xff, 0xff)

# Fonts. Single serif per page. PPT falls back on the viewer's system.
CN_SERIF = "LXGW WenKai"

SERIF = CN_SERIF
SANS  = SERIF

# 16:9 宽屏
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ═══════════════════════════════════════════════════════════
# Helpers
# ═══════════════════════════════════════════════════════════

def blank_slide(prs, bg_color=PARCHMENT):
    """创建空白幻灯片，指定背景色"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # 6 = Blank
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def add_text(slide, text, left, top, width, height,
             font=SANS, size=18, bold=False, italic=False,
             color=NEAR_BLACK, align=PP_ALIGN.LEFT,
             vanchor=MSO_ANCHOR.TOP):
    """加一段文字"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = vanchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_line(slide, left, top, width, color=BRAND, weight_pt=1):
    """加水平线"""
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    return line


def add_card(slide, left, top, width, height,
             fill=IVORY, border=BORDER, border_weight=0.5):
    """加卡片背景"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(border_weight)
    card.shadow.inherit = False
    return card


def add_diagram_png(slide, path, left, top, width, height):
    """Place a generated diagram PNG derived from the single-source SVG."""
    return slide.shapes.add_picture(str(path), left, top, width=width, height=height)


# ═══════════════════════════════════════════════════════════
# Slide Templates
# ═══════════════════════════════════════════════════════════

def cover_slide(prs, title, subtitle, author, date):
    """封面：大标题 + 副标题 + 作者/日期"""
    s = blank_slide(prs)
    # 大标题（serif 44pt 居中）
    add_text(s, title,
             Inches(1), Inches(2.5), Inches(11.33), Inches(1.5),
             font=SERIF, size=44, color=NEAR_BLACK,
             align=PP_ALIGN.CENTER)
    # 品牌色短线
    add_line(s, Inches(6.17), Inches(4.2), Inches(1), weight_pt=1.5)
    # 副标题
    add_text(s, subtitle,
             Inches(1), Inches(4.5), Inches(11.33), Inches(0.8),
             font=SANS, size=18, color=OLIVE,
             align=PP_ALIGN.CENTER)
    # 作者 + 日期
    add_text(s, f"{author}　·　{date}",
             Inches(1), Inches(6.5), Inches(11.33), Inches(0.4),
             font=SANS, size=13, color=STONE,
             align=PP_ALIGN.CENTER)
    return s


def toc_slide(prs, items):
    """目录页：01 章节名 列表"""
    s = blank_slide(prs)
    add_text(s, "目录",
             Inches(1.2), Inches(0.8), Inches(10), Inches(0.8),
             font=SERIF, size=32, color=NEAR_BLACK)
    add_line(s, Inches(1.2), Inches(1.8), Inches(11), weight_pt=1)

    for i, item in enumerate(items):
        y = Inches(2.4 + i * 0.9)
        add_text(s, f"0{i+1}",
                 Inches(1.2), y, Inches(1), Inches(0.6),
                 font=SERIF, size=28, color=BRAND)
        add_text(s, item,
                 Inches(2.4), y, Inches(9), Inches(0.6),
                 font=SERIF, size=22, color=NEAR_BLACK,
                 vanchor=MSO_ANCHOR.MIDDLE)
    return s


def chapter_slide(prs, number, title):
    """章节首页：印章朱红色背景 + 居中大标题"""
    s = blank_slide(prs, bg_color=BRAND)
    add_text(s, f"0{number}",
             Inches(0.8), Inches(0.5), Inches(2), Inches(0.8),
             font=SERIF, size=26, color=WHITE)
    add_text(s, title,
             Inches(1), Inches(3), Inches(11.33), Inches(1.5),
             font=SERIF, size=56, color=WHITE,
             align=PP_ALIGN.CENTER)
    return s


def content_slide(prs, eyebrow, title, body, page_num=None):
    """内容页：小标题 + 大标题 + 正文"""
    s = blank_slide(prs)
    # eyebrow
    add_text(s, eyebrow,
             Inches(1.2), Inches(0.6), Inches(10), Inches(0.4),
             font=SANS, size=12, color=STONE)
    # title
    add_text(s, title,
             Inches(1.2), Inches(1.2), Inches(11.33), Inches(1.2),
             font=SERIF, size=32, color=NEAR_BLACK)
    # body
    add_text(s, body,
             Inches(1.2), Inches(3), Inches(11), Inches(3.5),
             font=SANS, size=18, color=DARK_WARM)
    # page number
    if page_num is not None:
        add_text(s, f" - {page_num:02d}",
                 Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.3),
                 font=SANS, size=11, color=STONE,
                 align=PP_ALIGN.RIGHT)
    return s


def metrics_slide(prs, title, metrics):
    """数据页：标题 + N 张数据卡并排
    metrics: [(value, label), ...]
    """
    s = blank_slide(prs)
    # 标题
    add_text(s, title,
             Inches(1.2), Inches(0.8), Inches(11), Inches(1),
             font=SERIF, size=28, color=NEAR_BLACK,
             align=PP_ALIGN.CENTER)
    add_line(s, Inches(6.17), Inches(2), Inches(1))

    # 数据卡
    n = len(metrics)
    card_w = Inches(2.8)
    gap = Inches(0.3)
    total_w = card_w * n + gap * (n - 1)
    start = (SLIDE_W - total_w) / 2

    for i, (value, label) in enumerate(metrics):
        x = start + (card_w + gap) * i
        # 大数字
        add_text(s, value,
                 x, Inches(3), card_w, Inches(1.5),
                 font=SERIF, size=52, color=BRAND,
                 align=PP_ALIGN.CENTER)
        # 标签
        add_text(s, label,
                 x, Inches(4.8), card_w, Inches(0.6),
                 font=SANS, size=14, color=OLIVE,
                 align=PP_ALIGN.CENTER)
    return s


def quote_slide(prs, quote, source):
    """引用页：极简，居中引文"""
    s = blank_slide(prs)
    add_text(s, f"\u201c{quote}\u201d",
             Inches(1.5), Inches(2.8), Inches(10.33), Inches(2.5),
             font=SERIF, size=28, color=NEAR_BLACK,
             align=PP_ALIGN.CENTER,
             vanchor=MSO_ANCHOR.MIDDLE)
    add_text(s, f" - {source}",
             Inches(1.5), Inches(5.2), Inches(10.33), Inches(0.4),
             font=SANS, size=14, color=OLIVE,
             align=PP_ALIGN.CENTER)
    return s


def comparison_slide(prs, eyebrow, left_title, left_items, right_title, right_items, page_num=None):
    """对比页：左右两栏，竖线分隔，左侧降调，右侧全色
    left_items / right_items: list of str (最多 4 条)
    """
    s = blank_slide(prs)
    add_text(s, eyebrow,
             Inches(1.2), Inches(0.6), Inches(10), Inches(0.4),
             font=SANS, size=12, color=STONE)
    # 分隔竖线（居中）
    divider = s.shapes.add_connector(1,
        Inches(6.67), Inches(1.0),
        Inches(6.67), Inches(6.8))
    divider.line.color.rgb = BORDER
    divider.line.width = Pt(1)
    # 左栏标题（降调）
    add_text(s, left_title,
             Inches(1.2), Inches(1.2), Inches(5), Inches(0.8),
             font=SERIF, size=22, color=OLIVE)
    # 右栏标题（全色）
    add_text(s, right_title,
             Inches(7.0), Inches(1.2), Inches(5), Inches(0.8),
             font=SERIF, size=22, color=NEAR_BLACK)
    # 分隔线
    add_line(s, Inches(1.2), Inches(2.2), Inches(11.5), weight_pt=0.5)
    # 左栏条目（降调）
    for i, item in enumerate(left_items[:4]):
        add_text(s, item,
                 Inches(1.2), Inches(2.6 + i * 0.9), Inches(4.9), Inches(0.7),
                 font=SANS, size=17, color=STONE)
    # 右栏条目（全色）
    for i, item in enumerate(right_items[:4]):
        add_text(s, item,
                 Inches(7.0), Inches(2.6 + i * 0.9), Inches(5.2), Inches(0.7),
                 font=SANS, size=17, color=DARK_WARM)
    if page_num is not None:
        add_text(s, f" - {page_num:02d}",
                 Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.3),
                 font=SANS, size=11, color=STONE,
                 align=PP_ALIGN.RIGHT)
    return s


def pipeline_slide(prs, eyebrow, title, steps, page_num=None):
    """流程步骤页：01/02/03 序号 + 步骤标题 + 步骤描述
    steps: list of (step_title, step_desc)，最多 4 步
    """
    s = blank_slide(prs)
    add_text(s, eyebrow,
             Inches(1.2), Inches(0.6), Inches(10), Inches(0.4),
             font=SANS, size=12, color=STONE)
    add_text(s, title,
             Inches(1.2), Inches(1.1), Inches(11), Inches(0.9),
             font=SERIF, size=30, color=NEAR_BLACK)
    add_line(s, Inches(1.2), Inches(2.15), Inches(11), weight_pt=0.5)

    n = len(steps[:4])
    step_w = Inches(11.5 / n)
    for i, (step_title, step_desc) in enumerate(steps[:4]):
        x = Inches(1.0) + step_w * i
        # 序号
        add_text(s, f"0{i+1}",
                 x, Inches(2.5), step_w, Inches(0.8),
                 font=SERIF, size=40, color=BRAND)
        # 步骤标题
        add_text(s, step_title,
                 x, Inches(3.45), step_w - Inches(0.2), Inches(0.6),
                 font=SERIF, size=19, color=NEAR_BLACK)
        # 步骤描述
        add_text(s, step_desc,
                 x, Inches(4.15), step_w - Inches(0.2), Inches(2.2),
                 font=SANS, size=15, color=OLIVE)
    if page_num is not None:
        add_text(s, f" - {page_num:02d}",
                 Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.3),
                 font=SANS, size=11, color=STONE,
                 align=PP_ALIGN.RIGHT)
    return s


def ending_slide(prs, message, contact):
    """结束页"""
    s = blank_slide(prs)
    add_text(s, message,
             Inches(1), Inches(3), Inches(11.33), Inches(1.2),
             font=SERIF, size=40, color=NEAR_BLACK,
             align=PP_ALIGN.CENTER)
    add_line(s, Inches(6.17), Inches(4.5), Inches(1), weight_pt=1.5)
    add_text(s, contact,
             Inches(1), Inches(4.8), Inches(11.33), Inches(0.6),
             font=SANS, size=16, color=OLIVE,
             align=PP_ALIGN.CENTER)
    return s


# ═══════════════════════════════════════════════════════════
# Main: 示例 deck，按实际需求改
# ═══════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. 封面
    cover_slide(prs,
        title="{{文档标题}}",
        subtitle="{{一句话描述}}",
        author="{{作者}}",
        date="2026.04")

    # 2. 目录
    toc_slide(prs, items=[
        "{{章节 1}}",
        "{{章节 2}}",
        "{{章节 3}}",
        "{{Q&A}}",
    ])

    # 3. 章节首页
    chapter_slide(prs, 1, "{{章节标题}}")

    # 5. 内容页
    content_slide(prs,
        eyebrow="{{章节 · 本页}}",
        title="{{核心论点标题}}",
        body="{{一段正文，18pt sans 字体。控制在 3 行内，一屏一个核心信息。}}",
        page_num=5)

    # 7. 数据页
    metrics_slide(prs,
        title="关键结果",
        metrics=[
            ("+42%",   "转化率提升"),
            ("3.8M",   "月活用户"),
            ("99.9%",  "可用性 SLA"),
            ("5,000+", "QPS 峰值"),
        ])

    # 6. 引用
    quote_slide(prs,
        quote="好的设计是尽可能少的设计。",
        source="Dieter Rams")

    # 8. 对比页
    comparison_slide(prs,
        eyebrow="{{章节 · 对比}}",
        left_title="{{旧方案}}",
        left_items=["{{对比点 A}}", "{{对比点 B}}", "{{对比点 C}}"],
        right_title="{{新方案}}",
        right_items=["{{改善点 A}}", "{{改善点 B}}", "{{改善点 C}}"],
        page_num=8)

    # 9. 流程步骤页
    pipeline_slide(prs,
        eyebrow="{{章节 · 流程}}",
        title="{{核心流程标题}}",
        steps=[
            ("{{步骤 1}}", "{{步骤 1 的说明文字，控制在两行内。}}"),
            ("{{步骤 2}}", "{{步骤 2 的说明文字，控制在两行内。}}"),
            ("{{步骤 3}}", "{{步骤 3 的说明文字，控制在两行内。}}"),
        ],
        page_num=9)

    # 7. 结束
    ending_slide(prs,
        message="Thank you",
        contact="{{邮箱}}　·　{{网站}}")

    prs.save('output.pptx')
    print("OK: Saved output.pptx")


if __name__ == '__main__':
    main()
