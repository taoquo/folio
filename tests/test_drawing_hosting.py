import json
import re
import sys
from pathlib import Path
from tempfile import TemporaryDirectory
from unittest import TestCase, mock

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))

from drawing.compiler import DEFAULT_COMPILER_REGISTRY
from drawing.hosting import (
    add_fitted_picture,
    embed_html_figure,
    embed_pptx_slot,
    fit_contain,
    host_contract,
    list_host_contracts,
    verify_hosted_html,
    verify_hosted_pptx,
)


def compile_minimal(kind: str):
    path = ROOT / "references" / "fixtures" / "minimal" / f"{kind}.json"
    payload = json.loads(path.read_text(encoding="utf-8"))
    return path, DEFAULT_COMPILER_REGISTRY.compile_payload(payload)


def compile_minimal_themed(kind: str, theme: str):
    path = ROOT / "references" / "fixtures" / "minimal" / f"{kind}.json"
    payload = json.loads(path.read_text(encoding="utf-8"))
    return path, DEFAULT_COMPILER_REGISTRY.compile_payload(payload, "embed", theme)


class DrawingHostingTests(TestCase):
    def test_four_explicit_host_contracts_have_bounded_safe_areas(self) -> None:
        contracts = list_host_contracts()

        self.assertEqual(
            {"a4-portrait", "letter-portrait", "slide-16x9", "responsive-html"},
            {item.key for item in contracts},
        )
        for item in contracts:
            with self.subTest(host=item.key):
                self.assertGreater(item.safe_width, 0)
                self.assertGreater(item.safe_height, 0)
                self.assertIn(item.default_profile, item.allowed_profiles)
        self.assertEqual("svg", host_contract("a4-portrait").artifact_format)
        self.assertEqual("png", host_contract("slide-16x9").artifact_format)

    def test_contain_fit_preserves_aspect_ratio_and_centers(self) -> None:
        left, top, width, height = fit_contain(400, 100, 10, 20, 200, 200)

        self.assertAlmostEqual(4, width / height)
        self.assertAlmostEqual(10, left)
        self.assertAlmostEqual(95, top)
        self.assertLessEqual(width, 200)
        self.assertLessEqual(height, 200)

    def test_html_embedding_supports_structural_and_exact_chart_fallback(self) -> None:
        tree_path, tree = compile_minimal("tree")
        bar_path, bar = compile_minimal("bar-chart")
        host = """<!DOCTYPE html><html lang="en"><head><meta charset="UTF-8"><style>@page { size: A4; }</style></head><body>
<figure data-folio-diagram-slot="structure"><p>replace</p></figure>
<figure data-folio-diagram-slot="data"><p>replace</p></figure>
</body></html>"""

        with TemporaryDirectory() as temp:
            root = Path(temp)
            source, output, artifacts = root / "host.html", root / "output.html", root / "artifacts"
            source.write_text(host, encoding="utf-8")
            embed_html_figure(
                tree, fixture=tree_path, host_file=source, output_host=output,
                artifact_dir=artifacts, contract_key="a4-portrait", slot="structure",
                caption="The root keeps this hierarchy intentionally bounded and easy to scan.",
            )
            embed_html_figure(
                bar, fixture=bar_path, host_file=output, output_host=output,
                artifact_dir=artifacts, contract_key="a4-portrait", slot="data",
                caption="The single category establishes a traceable baseline for later comparisons.",
            )
            manifests = verify_hosted_html(output)
            generated = output.read_text(encoding="utf-8")

            self.assertEqual(2, len(manifests))
            self.assertEqual(1, generated.count('id="folio-diagram-host-styles"'))
            self.assertIn('data-folio-data-table="true"', generated)
            self.assertIn('<th scope="col">Category</th>', generated)
            self.assertIn('<td>Category</td>', generated)
            self.assertIn('<td>1</td>', generated)
            self.assertIn("max-height:", generated)
            ids = re.findall(r'(?<![-\w])id="([^"]+)"', generated)
            self.assertEqual(len(ids), len(set(ids)))
            for labelled in re.findall(r'aria-labelledby="([^"]+)"', generated):
                self.assertTrue(set(labelled.split()) <= set(ids))

            chart = next(item for item in manifests if item["kind"] == "bar-chart")
            artifact = output.parent / chart["artifact"]
            original = output.read_text(encoding="utf-8")
            output.write_text(original.replace("Root", "Changed root", 1), encoding="utf-8")
            with self.assertRaisesRegex(ValueError, "embedded SVG is stale"):
                verify_hosted_html(output)
            output.write_text(original, encoding="utf-8")
            output.write_text(original.replace("<td>1</td>", "<td>2</td>", 1), encoding="utf-8")
            with self.assertRaisesRegex(ValueError, "accessible data table is stale"):
                verify_hosted_html(output)
            output.write_text(original, encoding="utf-8")
            artifact.write_text("stale", encoding="utf-8")
            with self.assertRaisesRegex(ValueError, "artifact is stale"):
                verify_hosted_html(output)

    def test_html_embedding_verifies_before_atomic_replacement(self) -> None:
        tree_path, tree = compile_minimal("tree")
        host = '<!DOCTYPE html><html><head></head><body><figure data-folio-diagram-slot="main"></figure></body></html>'
        with TemporaryDirectory() as temp:
            root = Path(temp)
            source, output = root / "source.html", root / "output.html"
            source.write_text(host, encoding="utf-8")
            output.write_text("previous valid output", encoding="utf-8")
            with mock.patch("drawing.hosting.verify_hosted_html", side_effect=ValueError("verification failed")):
                with self.assertRaisesRegex(ValueError, "verification failed"):
                    embed_html_figure(
                        tree, fixture=tree_path, host_file=source, output_host=output,
                        artifact_dir=root / "assets", contract_key="responsive-html", slot="main",
                        caption="The single root remains bounded and clearly readable in the responsive host.",
                    )
            self.assertEqual("previous valid output", output.read_text(encoding="utf-8"))

    def test_letter_and_responsive_html_contracts_embed_without_static_min_width(self) -> None:
        tree_path, tree = compile_minimal("tree")
        host = "<!DOCTYPE html><html><head></head><body><figure data-folio-diagram-slot=\"main\"></figure></body></html>"

        with TemporaryDirectory() as temp:
            root = Path(temp)
            for contract in ("letter-portrait", "responsive-html"):
                with self.subTest(contract=contract):
                    source = root / f"{contract}-source.html"
                    output = root / f"{contract}.html"
                    source.write_text(host, encoding="utf-8")
                    embed_html_figure(
                        tree, fixture=tree_path, host_file=source, output_host=output,
                        artifact_dir=root / f"{contract}-assets", contract_key=contract, slot="main",
                        caption="The bounded root remains legible across this host without a fixed minimum width.",
                    )
                    text = output.read_text(encoding="utf-8")
                    self.assertIn("min-width: 0 !important", text)
                    self.assertEqual(contract, verify_hosted_html(output)[0]["host_contract"])

    def test_pptx_slot_embedding_preserves_ratio_alt_text_notes_and_stale_gate(self) -> None:
        fixture, result = compile_minimal("line-chart")
        with TemporaryDirectory() as temp:
            root = Path(temp)
            source, output = root / "host.pptx", root / "output.pptx"
            prs = Presentation()
            prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(11), Inches(5.5))
            slot.name = "folio-diagram-slot:trend"
            prs.save(source)

            manifest = embed_pptx_slot(
                result, fixture=fixture, host_file=source, output_host=output,
                artifact_dir=root / "artifacts", slot="trend",
                caption="The second point rises above the first, establishing a clear upward direction.",
                slide_index=1, profile="embed",
            )
            manifests = verify_hosted_pptx(output)
            saved = Presentation(output)
            picture = next(shape for shape in saved.slides[0].shapes if shape.name.startswith("folio-diagram:trend:"))
            properties = picture._element.xpath(".//p:cNvPr")[0]

            self.assertEqual([manifest], manifests)
            self.assertAlmostEqual(960 / 540, picture.width / picture.height, delta=0.002)
            self.assertTrue(properties.get("descr"))
            notes = saved.slides[0].notes_slide.notes_text_frame.text
            self.assertIn("DATA SUMMARY:", notes)
            self.assertIn("Category\tSeries", notes)

            tampered = root / "tampered.pptx"
            tampered_prs = Presentation(output)
            frame = tampered_prs.slides[0].notes_slide.notes_text_frame
            frame.text = frame.text.replace("First\t1", "First\t9")
            tampered_prs.save(tampered)
            with self.assertRaisesRegex(ValueError, "accessible data notes are stale"):
                verify_hosted_pptx(tampered)

            artifact = output.parent / manifest["artifact"]
            artifact.write_bytes(b"stale")
            with self.assertRaisesRegex(ValueError, "artifact is stale"):
                verify_hosted_pptx(output)

    def test_slide_template_picture_helper_never_distorts(self) -> None:
        with TemporaryDirectory() as temp:
            image_path = Path(temp) / "wide.png"
            Image.new("RGB", (400, 100), "#F6F0EA").save(image_path)
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            picture = add_fitted_picture(
                slide, image_path, Inches(1), Inches(1), Inches(4), Inches(4),
                alt_text="A wide diagram", title="Wide",
            )

            self.assertAlmostEqual(4, picture.width / picture.height, delta=0.002)
            self.assertEqual(Inches(4), picture.width)
            self.assertEqual(Inches(1), picture.height)


class HostThemeVariantTests(TestCase):
    HOST = '<!DOCTYPE html><html lang="en"><head><meta charset="UTF-8"></head><body><figure data-folio-diagram-slot="main"></figure></body></html>'

    def test_html_manifest_records_theme_and_variant_and_svg_carries_them(self) -> None:
        fixture, result = compile_minimal_themed("tree", "dark")
        with TemporaryDirectory() as temp:
            root = Path(temp)
            source, output = root / "source.html", root / "output.html"
            source.write_text(self.HOST, encoding="utf-8")

            manifest = embed_html_figure(
                result, fixture=fixture, host_file=source, output_host=output,
                artifact_dir=root / "artifacts", contract_key="a4-portrait", slot="main",
                caption="The bounded root stays legible after the dark theme is applied to the scene.",
                variant="sketchy",
            )
            generated = output.read_text(encoding="utf-8")

            self.assertEqual("1.1", manifest["schema_version"])
            self.assertEqual("dark", manifest["theme"])
            self.assertEqual("sketchy", manifest["variant"])
            self.assertIn('data-folio-theme="dark"', generated)
            self.assertIn('data-folio-variant="sketchy"', generated)
            self.assertIn("feDisplacementMap", generated)
            self.assertEqual([manifest], verify_hosted_html(output))

    def test_html_variant_metadata_drift_is_rejected(self) -> None:
        fixture, result = compile_minimal("tree")
        with TemporaryDirectory() as temp:
            root = Path(temp)
            source, output = root / "source.html", root / "output.html"
            source.write_text(self.HOST, encoding="utf-8")
            embed_html_figure(
                result, fixture=fixture, host_file=source, output_host=output,
                artifact_dir=root / "artifacts", contract_key="a4-portrait", slot="main",
                caption="The bounded root remains readable while the variant metadata is checked for drift.",
            )
            original = output.read_text(encoding="utf-8")
            drifted = original.replace("&quot;variant&quot;:&quot;plain&quot;", "&quot;variant&quot;:&quot;sketchy&quot;", 1)
            self.assertNotEqual(original, drifted)
            output.write_text(drifted, encoding="utf-8")

            with self.assertRaisesRegex(ValueError, "variant metadata is stale"):
                verify_hosted_html(output)

    def test_unknown_variant_and_motion_pptx_are_rejected(self) -> None:
        fixture, result = compile_minimal("tree")
        with TemporaryDirectory() as temp:
            root = Path(temp)
            source, output = root / "source.html", root / "output.html"
            source.write_text(self.HOST, encoding="utf-8")

            with self.assertRaisesRegex(ValueError, "unknown drawing output variant"):
                embed_html_figure(
                    result, fixture=fixture, host_file=source, output_host=output,
                    artifact_dir=root / "artifacts", contract_key="a4-portrait", slot="main",
                    caption="The bounded root is never rendered because the requested variant does not exist.",
                    variant="glow",
                )

            deck = root / "deck.pptx"
            prs = Presentation()
            prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(11), Inches(5.5))
            slot.name = "folio-diagram-slot:main"
            prs.save(deck)

            with self.assertRaisesRegex(ValueError, "motion is CSS-driven"):
                embed_pptx_slot(
                    result, fixture=fixture, host_file=deck, output_host=root / "out.pptx",
                    artifact_dir=root / "artifacts", slot="main",
                    caption="The bounded root is never rasterised because motion cannot survive a static export.",
                    slide_index=1, profile="embed", variant="motion",
                )
