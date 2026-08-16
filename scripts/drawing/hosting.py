from __future__ import annotations

from dataclasses import asdict, dataclass
from hashlib import sha256
from html import escape
from html.parser import HTMLParser
import json
import os
from pathlib import Path
import re
from tempfile import NamedTemporaryFile, TemporaryDirectory
from typing import Any

from .compiler import CompilationResult
from .output import normalize_output_variant


MANIFEST_PREFIX = "FOLIO_DIAGRAM_MANIFEST:"
DATA_KINDS = {"bar-chart", "line-chart", "donut-chart", "candlestick", "waterfall"}


@dataclass(frozen=True)
class HostContract:
    key: str
    medium: str
    width: float
    height: float
    unit: str
    safe_left: float
    safe_top: float
    safe_right: float
    safe_bottom: float
    caption_gap: float
    artifact_format: str
    default_profile: str
    allowed_profiles: tuple[str, ...]
    target_ppi: int | None = None

    @property
    def safe_width(self) -> float:
        return self.width - self.safe_left - self.safe_right

    @property
    def safe_height(self) -> float:
        return self.height - self.safe_top - self.safe_bottom

    def to_dict(self) -> dict[str, Any]:
        payload = asdict(self)
        payload["allowed_profiles"] = list(self.allowed_profiles)
        payload["safe_width"] = self.safe_width
        payload["safe_height"] = self.safe_height
        return payload


HOST_CONTRACTS = {
    item.key: item
    for item in (
        HostContract("a4-portrait", "html-print", 210, 297, "mm", 22, 20, 22, 22, 2.1, "svg", "embed", ("embed",)),
        HostContract("letter-portrait", "html-print", 215.9, 279.4, "mm", 22, 20, 22, 22, 2.1, "svg", "embed", ("embed",)),
        HostContract("slide-16x9", "pptx", 13.333, 7.5, "in", 0.8, 0.6, 0.8, 0.6, 0.12, "png", "artifact", ("artifact", "embed"), 144),
        HostContract("responsive-html", "html-responsive", 960, 540, "px", 0, 0, 0, 0, 12, "svg", "embed", ("embed",)),
    )
}


@dataclass(frozen=True)
class AccessibleData:
    summary: str
    headers: tuple[str, ...]
    rows: tuple[tuple[str, ...], ...]

    def to_dict(self) -> dict[str, Any]:
        return {"summary": self.summary, "headers": list(self.headers), "rows": [list(row) for row in self.rows]}


def host_contract(key: str) -> HostContract:
    try:
        return HOST_CONTRACTS[key]
    except KeyError as exc:
        raise ValueError(f"unknown diagram host contract: {key}") from exc


def list_host_contracts() -> tuple[HostContract, ...]:
    return tuple(HOST_CONTRACTS[key] for key in sorted(HOST_CONTRACTS))


def fit_contain(
    source_width: float,
    source_height: float,
    box_left: float,
    box_top: float,
    box_width: float,
    box_height: float,
) -> tuple[float, float, float, float]:
    if min(source_width, source_height, box_width, box_height) <= 0:
        raise ValueError("source and host slot dimensions must be positive")
    scale = min(box_width / source_width, box_height / source_height)
    width, height = source_width * scale, source_height * scale
    left = box_left + (box_width - width) / 2
    top = box_top + (box_height - height) / 2
    return left, top, width, height


def validate_caption(caption: str, title: str) -> str:
    value = caption.strip()
    if len(value) < 16:
        raise ValueError("diagram caption must state an insight in at least 16 characters")
    if value.casefold().strip(" .。") == title.casefold().strip(" .。"):
        raise ValueError("diagram caption must state an insight, not repeat the title")
    blocked = ("chart title", "diagram title", "untitled", "图表标题", "图标题", "图表占位")
    if any(token in value.casefold() for token in blocked):
        raise ValueError("diagram caption contains placeholder language")
    return value


def accessible_data(result: CompilationResult) -> AccessibleData | None:
    kind = result.kind
    if kind not in DATA_KINDS:
        return None
    if kind in {"bar-chart", "line-chart"}:
        categories = result.plan.categories
        series = result.plan.series
        headers = ("Category", *(str(item["label"]) for item in series))
        rows = tuple(
            (str(category), *(_scalar(item["values"][index]) for item in series))
            for index, category in enumerate(categories)
        )
        summary = f"{len(categories)} categories across {len(series)} series."
    elif kind == "donut-chart":
        segments = result.plan.segments
        headers = ("Segment", "Value")
        rows = tuple((str(item["label"]), _scalar(item["value"])) for item in segments)
        total = sum(float(item["value"]) for item in segments)
        summary = f"{len(segments)} segments with total {_scalar(total)}."
    elif kind == "candlestick":
        periods = result.plan.periods
        headers = ("Date", "Open", "High", "Low", "Close")
        rows = tuple(
            tuple([str(item["date"]), *(_scalar(item[field]) for field in ("open", "high", "low", "close"))])
            for item in periods
        )
        summary = f"{len(periods)} periods from {periods[0]['date']} to {periods[-1]['date']}."
    else:
        start = float(result.plan.start)
        running = start
        rows_list = [("Start", "", _scalar(start))]
        for item in result.plan.contributions:
            if item.get("kind", "delta") == "subtotal":
                rows_list.append((str(item["label"]), "", _scalar(running)))
            else:
                running += float(item["value"])
                rows_list.append((str(item["label"]), _scalar(item["value"]), _scalar(running)))
        end = float(result.plan.end)
        rows_list.append(("End", "", _scalar(end)))
        headers = ("Step", "Change", "Running total")
        rows = tuple(rows_list)
        summary = f"Start {_scalar(start)}, {len(result.plan.contributions)} contributions, end {_scalar(end)}."
    return AccessibleData(summary, tuple(headers), tuple(tuple(value for value in row) for row in rows))


def _scalar(value: Any) -> str:
    if value is None:
        return "Missing"
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value)


def _digest(path: Path) -> str:
    return sha256(path.read_bytes()).hexdigest()


def _digest_text(value: str) -> str:
    return sha256(value.encode("utf-8")).hexdigest()


def _portable_path(path: Path, host_output: Path) -> str:
    return Path(os.path.relpath(path.resolve(), host_output.resolve().parent)).as_posix()


def _resolve_host_path(value: str, host_file: Path) -> Path:
    path = Path(value)
    return path if path.is_absolute() else host_file.parent / path


def _atomic_text(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with NamedTemporaryFile("w", encoding="utf-8", dir=path.parent, delete=False) as handle:
        handle.write(content)
        temporary = Path(handle.name)
    temporary.replace(path)


def _manifest(
    result: CompilationResult,
    fixture: Path,
    artifact: Path,
    host_output: Path,
    contract: HostContract,
    profile: str,
    slot: str,
    caption: str,
    data: AccessibleData | None,
    placement: dict[str, Any] | None = None,
    variant: str = "plain",
) -> dict[str, Any]:
    return {
        "schema_version": "1.1",
        "slot": slot,
        "host_contract": contract.key,
        "kind": result.kind,
        "profile": profile,
        "theme": result.theme,
        "variant": variant,
        "fixture": _portable_path(fixture, host_output),
        "fixture_sha256": _digest(fixture),
        "artifact": _portable_path(artifact, host_output),
        "artifact_sha256": _digest(artifact),
        "source_dimensions": [result.scene.width, result.scene.height],
        "caption": caption,
        "description": result.scene.description,
        "language": result.scene.language,
        "registry_key": result.metadata.registry_key if result.metadata else None,
        "data": data.to_dict() if data else None,
        "placement": placement,
    }


def embed_html_figure(
    result: CompilationResult,
    *,
    fixture: str | Path,
    host_file: str | Path,
    output_host: str | Path,
    artifact_dir: str | Path,
    contract_key: str,
    slot: str,
    caption: str,
    profile: str | None = None,
    variant: str = "plain",
) -> dict[str, Any]:
    from renderers.svg import render_svg

    contract = host_contract(contract_key)
    if contract.medium not in {"html-print", "html-responsive"}:
        raise ValueError(f"{contract.key} is not an HTML host contract")
    profile = profile or contract.default_profile
    if profile not in contract.allowed_profiles:
        raise ValueError(f"profile {profile} is not allowed for {contract.key}")
    variant = normalize_output_variant(variant)
    caption = validate_caption(caption, result.plan.title)
    fixture_path, source_path, output_path = Path(fixture), Path(host_file), Path(output_host)
    if not fixture_path.is_file() or not source_path.is_file():
        raise ValueError("fixture and HTML host file must exist")
    source = source_path.read_text(encoding="utf-8")
    marker = re.compile(
        rf"<figure\b(?=[^>]*\bdata-folio-diagram-slot=(['\"]){re.escape(slot)}\1)[^>]*>.*?</figure>",
        re.DOTALL,
    )
    matches = list(marker.finditer(source))
    if len(matches) != 1:
        raise ValueError(f"HTML host must contain exactly one figure slot named {slot}")

    artifact_root = Path(artifact_dir)
    artifact_root.mkdir(parents=True, exist_ok=True)
    token = _digest(fixture_path)[:10]
    svg = render_svg(result.scene, profile, namespace=f"host-{slot}-{token}", variant=variant)
    artifact = artifact_root / f"{result.kind}-{slot}-{token}-{_digest_text(svg)[:12]}.svg"
    _atomic_text(artifact, svg)
    data = accessible_data(result)
    max_height = contract.safe_height if contract.medium == "html-responsive" else round(contract.safe_height * 0.55, 3)
    placement = {
        "unit": contract.unit,
        "safe_area": [contract.safe_left, contract.safe_top, contract.safe_width, contract.safe_height],
        "max_artifact": [contract.safe_width, max_height],
    }
    manifest = _manifest(
        result, fixture_path, artifact, output_path, contract, profile, slot, caption, data, placement,
        variant=variant,
    )
    figure = _html_figure(result, svg, manifest, caption, data, contract)
    output = source[:matches[0].start()] + figure + source[matches[0].end():]
    output = _ensure_host_styles(output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with NamedTemporaryFile("w", encoding="utf-8", suffix=".html", dir=output_path.parent, delete=False) as handle:
        handle.write(output)
        temporary = Path(handle.name)
    try:
        verify_hosted_html(temporary)
        temporary.replace(output_path)
    finally:
        temporary.unlink(missing_ok=True)
    return manifest


def _html_figure(
    result: CompilationResult,
    svg: str,
    manifest: dict[str, Any],
    caption: str,
    data: AccessibleData | None,
    contract: HostContract,
) -> str:
    slot = str(manifest["slot"])
    caption_id = f"folio-caption-{slot}"
    table_id = f"folio-data-{slot}"
    attrs = {
        "data-folio-diagram-slot": slot,
        "data-folio-kind": str(manifest["kind"]),
        "data-folio-profile": str(manifest["profile"]),
        "data-folio-theme": str(manifest["theme"]),
        "data-folio-variant": str(manifest["variant"]),
        "data-folio-fixture-sha256": str(manifest["fixture_sha256"]),
        "data-folio-artifact-sha256": str(manifest["artifact_sha256"]),
        "role": "group",
        "aria-labelledby": caption_id,
    }
    if data:
        attrs["aria-describedby"] = table_id
    encoded_attrs = " ".join(f'{name}="{escape(value, quote=True)}"' for name, value in attrs.items())
    max_width, max_height = manifest["placement"]["max_artifact"]
    host_style = (
        f"--folio-host-max-width:{max_width}{contract.unit};"
        f"--folio-host-max-height:{max_height}{contract.unit}"
    )
    parts = [
        f"<figure {encoded_attrs}>",
        f'<div class="folio-diagram-artifact" style="{host_style}">', svg, "</div>",
    ]
    parts.append(f'<figcaption id="{caption_id}">{escape(caption)}</figcaption>')
    if data:
        parts.append(f'<div class="folio-data-fallback" id="{table_id}">')
        parts.append(f'<p class="folio-data-summary">{escape(data.summary)}</p>')
        parts.append('<table class="folio-table compact" data-folio-data-table="true"><thead><tr>')
        parts.extend(f'<th scope="col">{escape(header)}</th>' for header in data.headers)
        parts.append("</tr></thead><tbody>")
        for row in data.rows:
            parts.append("<tr>")
            parts.extend(f"<td>{escape(value)}</td>" for value in row)
            parts.append("</tr>")
        parts.append("</tbody></table></div>")
    manifest_json = json.dumps(manifest, ensure_ascii=False, sort_keys=True, separators=(",", ":"))
    parts.append(f'<script type="application/json" class="folio-diagram-manifest">{escape(manifest_json)}</script>')
    parts.append("</figure>")
    return "\n".join(parts)


HOST_STYLES = """<style id="folio-diagram-host-styles">
  .folio-diagram-artifact { display: flex; justify-content: center; width: 100%; max-width: var(--folio-host-max-width, 100%); overflow: hidden; }
  .folio-diagram-artifact svg { display: block; width: auto; max-width: 100%; height: auto; max-height: var(--folio-host-max-height, none); min-width: 0 !important; }
  .folio-data-fallback { margin-top: 8pt; break-inside: avoid; }
  .folio-data-summary { margin: 0 0 4pt; color: #5A4A43; font-size: 8.5pt; }
  .folio-data-fallback table { width: 100%; border-collapse: collapse; font-size: 7.5pt; }
  .folio-data-fallback th { text-align: left; color: #5A4A43; border-bottom: 0.75pt solid #E9DED4; padding: 2pt 4pt; }
  .folio-data-fallback td { border-bottom: 0.5pt solid #E9DED4; padding: 2pt 4pt; }
</style>"""


def _ensure_host_styles(source: str) -> str:
    if 'id="folio-diagram-host-styles"' in source:
        return source
    if "</head>" not in source:
        raise ValueError("HTML host must contain a head element")
    return source.replace("</head>", HOST_STYLES + "\n</head>", 1)


def embed_pptx_slot(
    result: CompilationResult,
    *,
    fixture: str | Path,
    host_file: str | Path,
    output_host: str | Path,
    artifact_dir: str | Path,
    slot: str,
    caption: str,
    slide_index: int,
    profile: str | None = None,
    variant: str = "plain",
) -> dict[str, Any]:
    from diagram_export import export_png
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    from renderers.svg import render_svg

    contract = host_contract("slide-16x9")
    profile = profile or contract.default_profile
    if profile not in contract.allowed_profiles:
        raise ValueError(f"profile {profile} is not allowed for {contract.key}")
    variant = normalize_output_variant(variant)
    if variant == "motion":
        raise ValueError("motion is CSS-driven and cannot be embedded into a raster PPTX slot")
    caption = validate_caption(caption, result.plan.title)
    fixture_path, source_path, output_path = Path(fixture), Path(host_file), Path(output_host)
    if not fixture_path.is_file() or not source_path.is_file():
        raise ValueError("fixture and PPTX host file must exist")
    prs = Presentation(str(source_path))
    if not 1 <= slide_index <= len(prs.slides):
        raise ValueError("slide_index is outside the presentation")
    slide = prs.slides[slide_index - 1]
    slot_name = f"folio-diagram-slot:{slot}"
    slots = [shape for shape in slide.shapes if shape.name == slot_name]
    if len(slots) != 1:
        raise ValueError(f"slide must contain exactly one image slot named {slot_name}")
    placeholder = slots[0]
    left, top, width, height = (placeholder.left, placeholder.top, placeholder.width, placeholder.height)
    _validate_slide_slot(left, top, width, height, contract)

    caption_height = Inches(0.36)
    caption_gap = Inches(contract.caption_gap)
    image_height = height - caption_height - caption_gap
    placement = fit_contain(
        result.scene.width, result.scene.height,
        float(left), float(top), float(width), float(image_height),
    )
    placeholder._element.getparent().remove(placeholder._element)

    artifact_root = Path(artifact_dir)
    artifact_root.mkdir(parents=True, exist_ok=True)
    token = _digest(fixture_path)[:10]
    target_width = int(round(contract.safe_width * int(contract.target_ppi or 144)))
    with TemporaryDirectory(prefix="folio-host-slide-", dir=artifact_root) as temp:
        svg_path = Path(temp) / "drawing.svg"
        svg_path.write_text(render_svg(result.scene, profile, variant=variant), encoding="utf-8")
        temporary_artifact = Path(temp) / "drawing.png"
        export_png(svg_path, temporary_artifact, width=target_width, profile=profile, title=result.plan.title, language=result.plan.language)
        artifact = artifact_root / f"{result.kind}-{slot}-{token}-{profile}-{_digest(temporary_artifact)[:12]}.png"
        temporary_artifact.replace(artifact)

    image = slide.shapes.add_picture(
        str(artifact), int(round(placement[0])), int(round(placement[1])),
        width=int(round(placement[2])), height=int(round(placement[3])),
    )
    image.name = f"folio-diagram:{slot}:{_digest(artifact)[:12]}"
    _set_picture_accessibility(image, result.scene.description, result.plan.title)

    caption_box = slide.shapes.add_textbox(left, top + image_height + caption_gap, width, caption_height)
    caption_box.name = f"folio-diagram-caption:{slot}"
    frame = caption_box.text_frame
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = caption
    run.font.name = "Charter"
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0x5A, 0x4A, 0x43)

    data = accessible_data(result)
    placement_manifest = {
        "unit": "in",
        "slot": [_emu_to_in(value) for value in (left, top, width, height)],
        "image": [_emu_to_in(value) for value in placement],
    }
    manifest = _manifest(
        result, fixture_path, artifact, output_path, contract, profile, slot, caption, data, placement_manifest,
        variant=variant,
    )
    notes = slide.notes_slide.notes_text_frame
    notes_lines = [notes.text.strip()] if notes.text.strip() else []
    notes_lines.append(MANIFEST_PREFIX + json.dumps(manifest, ensure_ascii=False, sort_keys=True, separators=(",", ":")))
    if data:
        notes_lines.append("DATA SUMMARY: " + data.summary)
        notes_lines.append("\t".join(data.headers))
        notes_lines.extend("\t".join(row) for row in data.rows)
    notes.text = "\n".join(notes_lines)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    with NamedTemporaryFile(suffix=".pptx", dir=output_path.parent, delete=False) as handle:
        temporary = Path(handle.name)
    try:
        prs.save(str(temporary))
        verify_hosted_pptx(temporary)
        temporary.replace(output_path)
    finally:
        temporary.unlink(missing_ok=True)
    return manifest


def _emu_to_in(value: float) -> float:
    return round(float(value) / 914400, 6)


def _validate_slide_slot(left: int, top: int, width: int, height: int, contract: HostContract) -> None:
    values = tuple(_emu_to_in(value) for value in (left, top, width, height))
    x, y, w, h = values
    tolerance = 0.002
    if x < contract.safe_left - tolerance or y < contract.safe_top - tolerance:
        raise ValueError("slide image slot begins outside the host safe area")
    if x + w > contract.width - contract.safe_right + tolerance or y + h > contract.height - contract.safe_bottom + tolerance:
        raise ValueError("slide image slot ends outside the host safe area")


def _set_picture_accessibility(shape: Any, description: str, title: str) -> None:
    properties = shape._element.xpath(".//p:cNvPr")
    if not properties:
        raise ValueError("PowerPoint image has no non-visual properties")
    properties[0].set("title", title)
    properties[0].set("descr", description)


def add_fitted_picture(
    slide: Any,
    path: str | Path,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    alt_text: str,
    title: str = "Folio diagram",
) -> Any:
    from PIL import Image

    image_path = Path(path)
    with Image.open(image_path) as image:
        source_width, source_height = image.size
    placement = fit_contain(source_width, source_height, left, top, width, height)
    picture = slide.shapes.add_picture(
        str(image_path), int(round(placement[0])), int(round(placement[1])),
        width=int(round(placement[2])), height=int(round(placement[3])),
    )
    _set_picture_accessibility(picture, alt_text, title)
    return picture


def verify_hosted_html(path: str | Path) -> list[dict[str, Any]]:
    host_file = Path(path)
    source = host_file.read_text(encoding="utf-8")
    manifests = _html_manifests(source)
    if not manifests:
        raise ValueError("HTML host contains no Folio diagram manifest")
    for manifest in manifests:
        _verify_manifest_files(manifest, host_file)
        slot = re.escape(str(manifest["slot"]))
        caption_match = re.search(
            rf"<figcaption\b[^>]*id=['\"]folio-caption-{slot}['\"][^>]*>(.*?)</figcaption>",
            source,
            re.DOTALL,
        )
        if not caption_match:
            raise ValueError(f"HTML host diagram {manifest['slot']} has no associated caption")
        if _plain_text(caption_match.group(1)) != manifest["caption"]:
            raise ValueError(f"HTML host diagram {manifest['slot']} caption is stale")
        figure_match = re.search(
            rf"<figure\b(?=[^>]*data-folio-diagram-slot=['\"]{slot}['\"])[^>]*>(.*?)</figure>",
            source,
            re.DOTALL,
        )
        if not figure_match:
            raise ValueError(f"HTML host diagram {manifest['slot']} figure is missing")
        svg_match = re.search(
            r'<div\b[^>]*class="[^"]*folio-diagram-artifact[^"]*"[^>]*>\s*(<svg\b.*?</svg>)\s*</div>',
            figure_match.group(1),
            re.DOTALL,
        )
        if not svg_match or _digest_text(svg_match.group(1)) != manifest["artifact_sha256"]:
            raise ValueError(f"HTML host diagram {manifest['slot']} embedded SVG is stale")
        if f'lang="{escape(str(manifest["language"]), quote=True)}"' not in svg_match.group(1):
            raise ValueError(f"HTML host diagram {manifest['slot']} language metadata is stale")
        if escape(str(manifest["description"])) not in svg_match.group(1):
            raise ValueError(f"HTML host diagram {manifest['slot']} description metadata is stale")
        if f'data-folio-variant="{escape(str(manifest["variant"]), quote=True)}"' not in svg_match.group(1):
            raise ValueError(f"HTML host diagram {manifest['slot']} variant metadata is stale")
        data = manifest.get("data")
        if manifest["kind"] in DATA_KINDS:
            if not data or f'id="folio-data-{manifest["slot"]}"' not in source:
                raise ValueError(f"HTML chart {manifest['slot']} has no accessible data table")
            parser = _AccessibleTableParser(f"folio-data-{manifest['slot']}")
            parser.feed(source)
            if parser.summary != data["summary"]:
                raise ValueError(f"HTML chart {manifest['slot']} data summary is stale")
            if parser.headers != data["headers"] or parser.rows != data["rows"]:
                raise ValueError(f"HTML chart {manifest['slot']} accessible data table is stale")
    return manifests


def _plain_text(source: str) -> str:
    parser = _TextParser()
    parser.feed(source)
    return "".join(parser.values).strip()


class _TextParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.values: list[str] = []

    def handle_data(self, data: str) -> None:
        self.values.append(data)


class _AccessibleTableParser(HTMLParser):
    def __init__(self, target_id: str):
        super().__init__(convert_charrefs=True)
        self.target_id = target_id
        self.depth = 0
        self.in_target = False
        self.capture: str | None = None
        self.buffer: list[str] = []
        self.summary = ""
        self.headers: list[str] = []
        self.rows: list[list[str]] = []
        self.current_row: list[str] | None = None

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        attributes = dict(attrs)
        if attributes.get("id") == self.target_id:
            self.in_target, self.depth = True, 1
            return
        if not self.in_target:
            return
        self.depth += 1
        if tag == "p" and "folio-data-summary" in str(attributes.get("class", "")).split():
            self.capture, self.buffer = "summary", []
        elif tag == "th":
            if attributes.get("scope") != "col":
                raise ValueError("accessible data table header must use scope=col")
            self.capture, self.buffer = "header", []
        elif tag == "tr":
            self.current_row = []
        elif tag == "td":
            self.capture, self.buffer = "cell", []

    def handle_endtag(self, tag: str) -> None:
        if not self.in_target:
            return
        if self.capture == "summary" and tag == "p":
            self.summary = "".join(self.buffer).strip()
            self.capture = None
        elif self.capture == "header" and tag == "th":
            self.headers.append("".join(self.buffer).strip())
            self.capture = None
        elif self.capture == "cell" and tag == "td":
            if self.current_row is not None:
                self.current_row.append("".join(self.buffer).strip())
            self.capture = None
        elif tag == "tr" and self.current_row is not None:
            if self.current_row:
                self.rows.append(self.current_row)
            self.current_row = None
        self.depth -= 1
        if self.depth == 0:
            self.in_target = False

    def handle_data(self, data: str) -> None:
        if self.in_target and self.capture:
            self.buffer.append(data)


def _html_manifests(source: str) -> list[dict[str, Any]]:
    pattern = re.compile(
        r'<script\s+type="application/json"\s+class="folio-diagram-manifest">(.*?)</script>',
        re.DOTALL,
    )
    from html import unescape

    return [json.loads(unescape(match.group(1))) for match in pattern.finditer(source)]


def verify_hosted_pptx(path: str | Path) -> list[dict[str, Any]]:
    from pptx import Presentation

    host_file = Path(path)
    prs = Presentation(str(host_file))
    manifests: list[dict[str, Any]] = []
    manifest_notes: dict[str, str] = {}
    for slide in prs.slides:
        notes_text = slide.notes_slide.notes_text_frame.text
        for line in notes_text.splitlines():
            if line.startswith(MANIFEST_PREFIX):
                manifest = json.loads(line[len(MANIFEST_PREFIX):])
                manifests.append(manifest)
                manifest_notes[str(manifest["slot"])] = notes_text
    if not manifests:
        raise ValueError("PPTX host contains no Folio diagram manifest")
    for manifest in manifests:
        _verify_manifest_files(manifest, host_file)
        name = f"folio-diagram:{manifest['slot']}:{manifest['artifact_sha256'][:12]}"
        shapes = [shape for slide in prs.slides for shape in slide.shapes if shape.name == name]
        if len(shapes) != 1:
            raise ValueError(f"PPTX host diagram {manifest['slot']} image is missing or duplicated")
        shape = shapes[0]
        properties = shape._element.xpath(".//p:cNvPr")
        if not properties or properties[0].get("descr") != manifest.get("description"):
            raise ValueError(f"PPTX host diagram {manifest['slot']} alt text is missing or stale")
        if not hasattr(shape, "image") or sha256(shape.image.blob).hexdigest() != manifest["artifact_sha256"]:
            raise ValueError(f"PPTX host diagram {manifest['slot']} embedded image is stale")
        source_width, source_height = manifest["source_dimensions"]
        if abs((shape.width / shape.height) - (source_width / source_height)) > 0.002:
            raise ValueError(f"PPTX host diagram {manifest['slot']} is distorted")
        captions = [
            item for slide in prs.slides for item in slide.shapes
            if item.name == f"folio-diagram-caption:{manifest['slot']}"
        ]
        if len(captions) != 1:
            raise ValueError(f"PPTX host diagram {manifest['slot']} caption is missing or duplicated")
        if captions[0].text.strip() != manifest["caption"]:
            raise ValueError(f"PPTX host diagram {manifest['slot']} caption is stale")
        data = manifest.get("data")
        if data:
            notes = manifest_notes[str(manifest["slot"])]
            expected_lines = [
                "DATA SUMMARY: " + data["summary"],
                "\t".join(data["headers"]),
                *("\t".join(row) for row in data["rows"]),
            ]
            if any(line not in notes for line in expected_lines):
                raise ValueError(f"PPTX host diagram {manifest['slot']} accessible data notes are stale")
    return manifests


def _verify_manifest_files(manifest: dict[str, Any], host_file: Path) -> None:
    required = {
        "schema_version", "slot", "host_contract", "kind", "profile", "theme", "variant", "fixture",
        "fixture_sha256", "artifact", "artifact_sha256", "source_dimensions", "caption",
        "description", "language", "registry_key",
    }
    if not required <= set(manifest):
        raise ValueError("diagram host manifest is incomplete")
    fixture = _resolve_host_path(str(manifest["fixture"]), host_file)
    artifact = _resolve_host_path(str(manifest["artifact"]), host_file)
    for label, file, digest_key in (
        ("fixture", fixture, "fixture_sha256"),
        ("artifact", artifact, "artifact_sha256"),
    ):
        if not file.is_file():
            raise ValueError(f"hosted diagram {label} is missing")
        if _digest(file) != manifest[digest_key]:
            raise ValueError(f"hosted diagram {label} is stale")


def verify_host(path: str | Path) -> list[dict[str, Any]]:
    host_file = Path(path)
    suffix = host_file.suffix.lower()
    if suffix in {".html", ".htm"}:
        return verify_hosted_html(host_file)
    if suffix == ".pptx":
        return verify_hosted_pptx(host_file)
    raise ValueError("diagram host must be HTML or PPTX")
