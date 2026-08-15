#!/usr/bin/env python3
from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


PARCHMENT = RGBColor(0xF6, 0xF0, 0xEA)
NEAR_BLACK = RGBColor(0x19, 0x15, 0x14)
BORDER = RGBColor(0xE9, 0xDE, 0xD4)


def add_title(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(1), Inches(0.55), Inches(11.3), Inches(0.6))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.name = "Charter"
    run.font.size = Pt(26)
    run.font.color.rgb = NEAR_BLACK


def add_slot(slide, slot: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.35), Inches(11.3), Inches(5.35))
    shape.name = f"folio-diagram-slot:{slot}"
    shape.fill.background()
    shape.line.color.rgb = BORDER


def build(path: Path) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    titles = (
        "Drawing Hosts",
        "Contract",
        "Architecture",
        "Evidence",
        "Trend",
        "Verification",
        "Next Step",
    )
    for index, title in enumerate(titles, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = PARCHMENT
        add_title(slide, title)
        if index == 3:
            add_slot(slide, "architecture")
        if index == 5:
            add_slot(slide, "trend")
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("output", type=Path)
    args = parser.parse_args()
    build(args.output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
